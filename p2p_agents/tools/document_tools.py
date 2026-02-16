"""Document generation tools — create PPT, Spreadsheet, and Doc files."""

from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path


GENERATED_DOCS_DIR = Path("generated_docs")


def _ensure_output_dir() -> Path:
    """Create the output directory if it doesn't exist and return it."""
    GENERATED_DOCS_DIR.mkdir(parents=True, exist_ok=True)
    return GENERATED_DOCS_DIR


def _timestamped_name(base: str, ext: str) -> str:
    """Return a filename like 'base_20260216_143022.ext'."""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_base = base.replace(" ", "_").replace("/", "_")[:60]
    return f"{safe_base}_{ts}.{ext}"


# ---------------------------------------------------------------------------
# Tool 1: PowerPoint
# ---------------------------------------------------------------------------

def create_ppt_report(
    title: str,
    slides: list[dict],
    template_style: str = "default",
) -> dict:
    """Creates a PowerPoint presentation and saves it to disk.

    Args:
        title: Presentation title shown on the first slide.
        slides: List of slide dicts. Each dict may contain:
            - heading (str): slide title
            - content (str): newline-separated bullet text
            - slide_type (str): 'bullets' (default), 'table', or 'title_only'
            - table_data (list[list[str]]): rows for table slides
              (first row treated as header)
        template_style: 'default' (white bg, blue accent) or
            'dark' (dark bg, white text).

    Returns:
        dict with status, file_path, slide_count, and title.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    out_dir = _ensure_output_dir()
    fname = _timestamped_name(title, "pptx")
    fpath = out_dir / fname

    is_dark = template_style.lower() == "dark"
    bg_color = RGBColor(0x1A, 0x1A, 0x2E) if is_dark else RGBColor(0xFF, 0xFF, 0xFF)
    text_color = RGBColor(0xFF, 0xFF, 0xFF) if is_dark else RGBColor(0x1A, 0x1A, 0x2E)
    accent_color = RGBColor(0x00, 0x7B, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title slide ---
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = accent_color

    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"ShareChat Finance — P2P Operations"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = datetime.now().strftime("%B %d, %Y")
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
    p3.alignment = PP_ALIGN.CENTER

    # --- Content slides ---
    for s in slides:
        heading = s.get("heading", "")
        content = s.get("content", "")
        slide_type = s.get("slide_type", "bullets")
        table_data = s.get("table_data", [])

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Header bar
        header_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(1.0)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = accent_color
        header_shape.line.fill.background()
        htf = header_shape.text_frame
        hp = htf.paragraphs[0]
        hp.text = heading
        hp.font.size = Pt(28)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        hp.alignment = PP_ALIGN.LEFT
        htf.margin_left = Inches(0.5)

        if slide_type == "table" and table_data:
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 1
            table = slide.shapes.add_table(
                rows, cols, Inches(0.5), Inches(1.3),
                Inches(12), Inches(0.4 * rows)
            ).table

            for ri, row in enumerate(table_data):
                for ci, cell_val in enumerate(row):
                    cell = table.cell(ri, ci)
                    cell.text = str(cell_val)
                    for para in cell.text_frame.paragraphs:
                        para.font.size = Pt(12)
                        if ri == 0:
                            para.font.bold = True
                            para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    if ri == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = accent_color

        elif slide_type == "title_only":
            pass  # just the header bar

        else:  # bullets
            txBox = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            lines = content.split("\n") if content else []
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line.strip().lstrip("•-").strip()
                p.font.size = Pt(18)
                p.font.color.rgb = text_color
                p.space_before = Pt(8)
                p.level = 0

    prs.save(str(fpath))

    return {
        "status": "created",
        "file_path": str(fpath.resolve()),
        "slide_count": len(slides) + 1,
        "title": title,
    }


# ---------------------------------------------------------------------------
# Tool 2: Spreadsheet (Excel)
# ---------------------------------------------------------------------------

def create_spreadsheet(
    title: str,
    sheets: list[dict],
) -> dict:
    """Creates an Excel spreadsheet and saves it to disk.

    Args:
        title: Workbook title (used as the filename).
        sheets: List of sheet dicts. Each dict should contain:
            - name (str): sheet tab name
            - headers (list[str]): column header labels
            - rows (list[list]): data rows (each inner list is one row)
            - column_widths (list[int], optional): per-column widths

    Returns:
        dict with status, file_path, sheet_count, and total_rows.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    out_dir = _ensure_output_dir()
    fname = _timestamped_name(title, "xlsx")
    fpath = out_dir / fname

    wb = Workbook()
    # Remove default sheet — we'll create named ones
    wb.remove(wb.active)

    header_font = Font(bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="007BFF", end_color="007BFF", fill_type="solid")
    header_align = Alignment(horizontal="center", vertical="center")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )
    total_rows = 0

    for sheet_def in sheets:
        name = sheet_def.get("name", "Sheet")[:31]  # Excel 31-char limit
        headers = sheet_def.get("headers", [])
        rows = sheet_def.get("rows", [])
        col_widths = sheet_def.get("column_widths", [])

        ws = wb.create_sheet(title=name)

        # Write headers
        for ci, h in enumerate(headers, start=1):
            cell = ws.cell(row=1, column=ci, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_align
            cell.border = thin_border

        # Write data
        for ri, row in enumerate(rows, start=2):
            for ci, val in enumerate(row, start=1):
                cell = ws.cell(row=ri, column=ci, value=val)
                cell.border = thin_border
                cell.alignment = Alignment(horizontal="left")
            total_rows += 1

        # Column widths
        for ci, w in enumerate(col_widths, start=1):
            col_letter = ws.cell(row=1, column=ci).column_letter
            ws.column_dimensions[col_letter].width = w

        # Auto-size columns that don't have explicit widths
        for ci in range(len(col_widths) + 1, len(headers) + 1):
            max_len = len(str(headers[ci - 1])) if ci <= len(headers) else 10
            for ri in range(2, len(rows) + 2):
                val = ws.cell(row=ri, column=ci).value
                if val is not None:
                    max_len = max(max_len, len(str(val)))
            col_letter = ws.cell(row=1, column=ci).column_letter
            ws.column_dimensions[col_letter].width = min(max_len + 4, 50)

        # Freeze top row
        ws.freeze_panes = "A2"

    wb.save(str(fpath))

    return {
        "status": "created",
        "file_path": str(fpath.resolve()),
        "sheet_count": len(sheets),
        "total_rows": total_rows,
    }


# ---------------------------------------------------------------------------
# Tool 3: Word Document
# ---------------------------------------------------------------------------

def create_doc_report(
    title: str,
    sections: list[dict],
    header_text: str = "ShareChat Finance — P2P Operations",
) -> dict:
    """Creates a Word document report and saves it to disk.

    Args:
        title: Document title (rendered as Heading 1).
        sections: List of section dicts. Each dict may contain:
            - heading (str): section heading (Heading 2)
            - content (str): paragraph text
            - bullet_points (list[str]): optional bullet list items
            - table (dict): optional table with 'headers' (list[str])
              and 'rows' (list[list[str]])
        header_text: Text for the page header.

    Returns:
        dict with status, file_path, section_count, and page_estimate.
    """
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    out_dir = _ensure_output_dir()
    fname = _timestamped_name(title, "docx")
    fpath = out_dir / fname

    doc = Document()

    # Page header
    section = doc.sections[0]
    header = section.header
    hp = header.paragraphs[0]
    hp.text = header_text
    hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = hp.runs[0]
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # Footer with date
    footer = section.footer
    fp = footer.paragraphs[0]
    fp.text = f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}"
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    frun = fp.runs[0]
    frun.font.size = Pt(8)
    frun.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # Title
    title_para = doc.add_heading(title, level=0)

    # Sections
    for sec in sections:
        heading = sec.get("heading", "")
        content = sec.get("content", "")
        bullets = sec.get("bullet_points", [])
        table_def = sec.get("table", None)

        if heading:
            doc.add_heading(heading, level=2)

        if content:
            p = doc.add_paragraph(content)
            p.paragraph_format.space_after = Pt(6)

        for bullet in bullets:
            doc.add_paragraph(bullet, style="List Bullet")

        if table_def:
            headers = table_def.get("headers", [])
            rows = table_def.get("rows", [])
            if headers:
                table = doc.add_table(
                    rows=1 + len(rows), cols=len(headers), style="Light Grid Accent 1"
                )
                # Header row
                for ci, h in enumerate(headers):
                    table.rows[0].cells[ci].text = str(h)
                    for para in table.rows[0].cells[ci].paragraphs:
                        for run in para.runs:
                            run.bold = True
                # Data rows
                for ri, row in enumerate(rows):
                    for ci, val in enumerate(row):
                        if ci < len(headers):
                            table.rows[ri + 1].cells[ci].text = str(val)
                doc.add_paragraph()  # spacing after table

    doc.save(str(fpath))

    # Rough page estimate: ~40 lines per page
    line_count = sum(
        1 + len(s.get("content", "").split("\n"))
        + len(s.get("bullet_points", []))
        + len(s.get("table", {}).get("rows", []))
        for s in sections
    )
    page_estimate = max(1, (line_count + 10) // 40 + 1)

    return {
        "status": "created",
        "file_path": str(fpath.resolve()),
        "section_count": len(sections),
        "page_estimate": page_estimate,
    }
